import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Define the data
data = [
    ("Feature Completeness", "9.0 / 10", "All core modules are fully functional and well-integrated."),
    ("Code Quality & Structure", "7.5 / 10", "Solid foundation using Vanilla JS and native HTML/CSS."),
    ("Security Practices", "7.5 / 10", "Strong baseline security with bcrypt password hashing."),
    ("Architecture & Design", "8.0 / 10", "Clean Node.js and Express architecture perfectly suited for MVP."),
    ("UI/UX Design", "9.0 / 10", "Highly polished, responsive, and consistent aesthetic."),
    ("Innovation & Creativity", "8.5 / 10", "Smart, creative use of client-side storage for prototyping."),
    ("Documentation Readiness", "7.5 / 10", "Code is highly readable, neatly formatted, and logically structured."),
    ("OVERALL", "8.2 / 10", "Very Good — A highly functional, beautifully designed application.")
]

prs = Presentation()

# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Vivienne - Technical Rating Review"
subtitle.text = "Comprehensive evaluation across key software engineering dimensions."

# Slide 2: Table Slide
blank_slide_layout = prs.slide_layouts[5] # Title only
slide2 = prs.slides.add_slide(blank_slide_layout)
shapes = slide2.shapes
title_shape = shapes.title
title_shape.text = "Evaluation Scorecard"

rows = len(data) + 1
cols = 3
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9.0)
height = Inches(0.8)

table = shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(1.5)
table.columns[2].width = Inches(5.0)

# Write headers
headers = ["Evaluation Dimension", "Score", "Remarks"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x00, 0x33, 0x66)
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            run.font.size = Pt(14)

# Write data
for row_idx, row_data in enumerate(data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = str(val)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(12)
                if row_idx == len(data) - 1: # Make the OVERALL row bold
                    run.font.bold = True

prs.save('/Users/omkarhattikal/Desktop/Vivienne_Rating_Review.pptx')
print("Successfully generated /Users/omkarhattikal/Desktop/Vivienne_Rating_Review.pptx")
